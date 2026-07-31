from __future__ import annotations

import io
import re
from pathlib import Path
from typing import Any, Sequence, Tuple

import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor

from defense_core import frequency, pct, confidence

BLANK_BLITZ = {"-", "NONE", "NO", "NO BLITZ", "0", "BASE", "NO DATA", "UNKNOWN", ""}


def _norm_col(name: str) -> str:
    return re.sub(r"[^a-z0-9]+", "", str(name).strip().lower())


def _fmt_pct(x) -> str:
    try:
        return f"{float(x):.1f}%"
    except Exception:
        return str(x)


def _count_pct(n: int | float, d: int | float) -> str:
    n = int(n or 0)
    d = int(d or 0)
    return f"{n}/{d} ({pct(n, d):.1f}%)"


def _shape_text(shape) -> str:
    return getattr(shape, "text", "") if hasattr(shape, "text") else ""


def _tables(slide):
    return [s for s in slide.shapes if getattr(s, "has_table", False)]


def _text_shapes(slide):
    return [s for s in slide.shapes if hasattr(s, "text_frame")]


def _copy_font(src, dst) -> None:
    """Copy explicit font settings so data replacement does not restyle the template."""
    try:
        dst.name = src.name
        dst.size = src.size
        dst.bold = src.bold
        dst.italic = src.italic
        dst.underline = src.underline
        if src.color and src.color.type is not None and src.color.rgb is not None:
            dst.color.rgb = src.color.rgb
    except Exception:
        pass


def _first_run_style(text_frame):
    for p in text_frame.paragraphs:
        for r in p.runs:
            return r
    return None


def _set_text_preserve_textframe(text_frame, value) -> None:
    """
    Replace visible text while preserving template layout and explicit font style.
    This does not touch fills, borders, margins, row heights, column widths, positions, or theme formatting.
    """
    value = "" if value is None else str(value)
    src_run = _first_run_style(text_frame)
    first_para = text_frame.paragraphs[0] if text_frame.paragraphs else None
    alignment = getattr(first_para, "alignment", None) if first_para is not None else None
    level = getattr(first_para, "level", None) if first_para is not None else None

    text_frame.clear()
    lines = value.split("\n")
    for i, line in enumerate(lines if lines else [""]):
        p = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
        if alignment is not None:
            p.alignment = alignment
        if level is not None:
            p.level = level
        run = p.add_run()
        run.text = line
        if src_run is not None:
            _copy_font(src_run.font, run.font)


def _set_shape_text_preserve(shape, value) -> None:
    if hasattr(shape, "text_frame"):
        _set_text_preserve_textframe(shape.text_frame, value)


def _set_cell_text_preserve(cell, value) -> None:
    _set_text_preserve_textframe(cell.text_frame, value)


def _fill_table_preserve(table_shape, headers: Sequence[str], rows: Sequence[Sequence[object]]) -> None:
    tbl = table_shape.table
    max_rows = len(tbl.rows)
    max_cols = len(tbl.columns)

    for c in range(max_cols):
        _set_cell_text_preserve(tbl.cell(0, c), headers[c] if c < len(headers) else "")

    for r in range(1, max_rows):
        for c in range(max_cols):
            val = ""
            if r - 1 < len(rows) and c < len(rows[r - 1]):
                val = rows[r - 1][c]
            _set_cell_text_preserve(tbl.cell(r, c), val)


def _replace_opponent(slide, opponent: str) -> None:
    for shape in _text_shapes(slide):
        txt = _shape_text(shape)
        if "Opponent:" in txt:
            _set_shape_text_preserve(shape, f"Opponent: {opponent}")


def _write_takeaways(slide, lines: Sequence[str]) -> None:
    for shape in _text_shapes(slide):
        txt = _shape_text(shape)
        if "Key Takeaways" in txt:
            _set_shape_text_preserve(shape, "Key Takeaways\n" + "\n".join(lines[:5]))
            return


def _top_value(df: pd.DataFrame, col: str, denom: int | None = None) -> Tuple[str, int, int, float]:
    f = frequency(df, col, denom=denom)
    if f.empty:
        return "No data", 0, int(denom or len(df)), 0.0
    row = f.iloc[0]
    d = int(denom if denom is not None else len(df))
    return str(row[col]), int(row["Plays"]), d, float(row["Pct"])


def _top_text(df: pd.DataFrame, col: str, denom: int | None = None) -> str:
    label, n, d, _ = _top_value(df, col, denom=denom)
    return "No data" if n == 0 else f"{label} — {_count_pct(n, d)}"


def _top_fronts_for_formation(g: pd.DataFrame) -> str:
    fronts = frequency(g, "FRONT").head(2)
    if fronts.empty:
        return "No data"
    total = len(g)
    top_pct = float(fronts.iloc[0]["Pct"])
    rows = fronts.iloc[:1] if top_pct >= 50 or len(fronts) == 1 else fronts.iloc[:2]
    return "\n".join(f'{r["FRONT"]} — {_count_pct(int(r["Plays"]), total)}' for _, r in rows.iterrows())


def _top_combos_text(df: pd.DataFrame, n=5) -> str:
    f = frequency(df, "COMBO").head(n)
    total = len(df)
    if f.empty or total == 0:
        return "No data"
    return "\n".join(
        [f'{i+1}. {r["COMBO"]} — {_count_pct(int(r["Plays"]), total)}' for i, r in f.iterrows()]
    )


def _third_down_bucket(df):
    d = df[(df["DOWN"] == 3)].copy()
    return {
        "3rd & 7+": d[d["DISTANCE"] >= 7],
        "3rd & 3-6": d[(d["DISTANCE"] >= 3) & (d["DISTANCE"] <= 6)],
        "3rd & 1-2": d[(d["DISTANCE"] >= 1) & (d["DISTANCE"] <= 2)],
    }


def _red_zone_bucket(df):
    y = df["YARD_LINE"]
    return {
        "High Red Zone (25-15)": df[(y >= 15) & (y <= 25)],
        "Red Zone (15-5)": df[(y >= 5) & (y < 15)],
        "Goal Line (Inside the 5)": df[(y >= 1) & (y < 5)],
    }


def _situational_rows_for_template(buckets: dict[str, pd.DataFrame], table_shape) -> tuple[list[str], list[list[object]]]:
    cols = len(table_shape.table.columns)
    if cols >= 3:
        headers = ["Situation", "Total Plays", "Top 5 Front/Stunt/Blitz/Coverage Calls"]
        rows = [[name, len(g), _top_combos_text(g, 5)] for name, g in buckets.items()]
    else:
        headers = ["Situation", "Top 5 Front/Stunt/Blitz/Coverage Calls"]
        rows = [[name, _top_combos_text(g, 5)] for name, g in buckets.items()]
    return headers, rows


def _read_table_like(path: str | Path) -> list[pd.DataFrame]:
    path = Path(path)
    if not path.exists():
        return []
    if path.suffix.lower() in {".xlsx", ".xls"}:
        xls = pd.read_excel(path, sheet_name=None)
        return [v.dropna(how="all") for v in xls.values()]
    return [pd.read_csv(path).dropna(how="all")]


def read_opponent_overview(path: str | Path | None) -> dict[str, Any]:
    """
    Optional local opponent-overview file reader. This replaces the unreliable website fetch.
    Supports CSV or Excel. Flexible columns:
    Snapshot fields: Team/Opponent, Overall Record/Record, Sacks, Interceptions/INTs, Fumble Recoveries.
    Schedule fields: Game/Week, Opponent, Record/Opponent Record, Score, W/L/Result.
    """
    out: dict[str, Any] = {"stats": {}, "schedule": []}
    if not path:
        return out

    for df in _read_table_like(path):
        if df.empty:
            continue
        df = df.copy()
        norm = {_norm_col(c): c for c in df.columns}

        # Key/value style sheet: Statistic | Value
        if "statistic" in norm and "value" in norm:
            for _, row in df.iterrows():
                key = str(row[norm["statistic"]]).strip()
                val = "" if pd.isna(row[norm["value"]]) else str(row[norm["value"]]).strip()
                if key:
                    out["stats"][key] = val
            continue

        # Schedule/game rows
        game_col = next((norm[x] for x in ["game", "week", "date"] if x in norm), None)
        opp_col = next((norm[x] for x in ["opponent", "opp"] if x in norm), None)
        score_col = next((norm[x] for x in ["score", "finalscore"] if x in norm), None)
        wl_col = next((norm[x] for x in ["wl", "w/l", "result", "winloss"] if x in norm), None)
        rec_col = next((norm[x] for x in ["record", "opponentrecord", "opprecord"] if x in norm), None)
        is_schedule = bool(opp_col and (score_col or wl_col))
        if is_schedule:
            for _, row in df.iterrows():
                out["schedule"].append({
                    "Game": "" if game_col is None or pd.isna(row[game_col]) else str(row[game_col]).strip(),
                    "Opponent": "" if pd.isna(row[opp_col]) else str(row[opp_col]).strip(),
                    "Record": "" if rec_col is None or pd.isna(row[rec_col]) else str(row[rec_col]).strip(),
                    "Score": "" if score_col is None or pd.isna(row[score_col]) else str(row[score_col]).strip(),
                    "W/L": "" if wl_col is None or pd.isna(row[wl_col]) else str(row[wl_col]).strip(),
                })

        # Snapshot row style. Do not treat schedule rows as season-record rows.
        aliases = {
            "Overall Record": ["overallrecord", "seasonrecord"],
            "Total Sacks": ["totalsacks", "sacks", "sack"],
            "Total INTs": ["totalints", "ints", "interceptions", "interception"],
            "Total Fumble Recoveries": ["totalfumblerecoveries", "fumblerecoveries", "fumblesrecovered", "fr"],
        }
        if not is_schedule and any(a in norm for vals in aliases.values() for a in vals):
            row = df.iloc[0]
            for label, vals in aliases.items():
                for v in vals:
                    if v in norm:
                        val = row[norm[v]]
                        out["stats"][label] = "" if pd.isna(val) else str(val).strip()
                        break
    return out



def _color_wl_cells(table_shape) -> None:
    """Color only the W/L result cells on the overview schedule table."""
    try:
        tbl = table_shape.table
        if len(tbl.columns) < 5:
            return
        for r in range(1, len(tbl.rows)):
            txt = tbl.cell(r, 4).text.strip().upper()
            if not txt:
                continue
            cell = tbl.cell(r, 4)
            cell.fill.solid()
            if txt.startswith("W"):
                cell.fill.fore_color.rgb = RGBColor(198, 239, 206)
            elif txt.startswith("L"):
                cell.fill.fore_color.rgb = RGBColor(255, 199, 206)
    except Exception:
        pass

def _fill_overview_slide(slide, opponent: str, overview: dict[str, Any] | None) -> None:
    overview = overview or {"stats": {}, "schedule": []}
    for shape in _text_shapes(slide):
        txt = _shape_text(shape)
        if "Opponent:" in txt:
            _set_shape_text_preserve(shape, f"Opponent: {opponent}")
    tables = _tables(slide)
    if not tables:
        return

    stats = overview.get("stats", {}) or {}
    stat_rows = [
        ["Overall Record", stats.get("Overall Record", "")],
        ["Total Sacks", stats.get("Total Sacks", "")],
        ["Total INTs", stats.get("Total INTs", "")],
        ["Total Fumble Recoveries", stats.get("Total Fumble Recoveries", "")],
    ]
    _fill_table_preserve(tables[0], ["Statistic", "Value"], stat_rows)

    if len(tables) > 1:
        schedule = overview.get("schedule", []) or []
        rows = [[r.get("Game", ""), r.get("Opponent", ""), r.get("Record", ""), r.get("Score", ""), r.get("W/L", "")] for r in schedule[:12]]
        _fill_table_preserve(tables[1], ["Game", "Opponent", "Record", "Score", "W/L"], rows)
        _color_wl_cells(tables[1])


def _find_slide_index(prs: Presentation, text: str) -> int | None:
    needle = text.lower()
    for i, slide in enumerate(prs.slides):
        if any(needle in _shape_text(s).lower() for s in _text_shapes(slide)):
            return i
    return None


def build_defense_pptx(engine, template_path: str | Path, opponent: str = "Opponent", overview: dict[str, Any] | None = None) -> bytes:
    prs = Presentation(str(template_path))
    df = engine.df
    total = len(df)

    for slide in prs.slides:
        _replace_opponent(slide, opponent)

    idx = {
        "overview": _find_slide_index(prs, "Opponent Overview"),
        "front": _find_slide_index(prs, "Front Tendencies"),
        "blitz": _find_slide_index(prs, "Blitz Tendencies"),
        "coverage": _find_slide_index(prs, "Coverage Tendencies"),
        "third": _find_slide_index(prs, "3rd Down Tendencies"),
        "redzone": _find_slide_index(prs, "Red Zone Tendencies"),
        "formation": _find_slide_index(prs, "Formation Tendencies"),
    }

    # Optional Slide 1: Opponent Overview
    if idx["overview"] is not None:
        _fill_overview_slide(prs.slides[idx["overview"]], opponent, overview)

    # Fronts
    if idx["front"] is not None:
        slide = prs.slides[idx["front"]]
        front_freq = frequency(df, "FRONT").head(5)
        rows = []
        for _, r in front_freq.iterrows():
            front = r["FRONT"]
            g = df[df["FRONT"] == front]
            front_snaps = len(g)
            blitz_count = int(g["IS_BLITZ"].sum())
            blitz_g = g[g["IS_BLITZ"]]
            top_blitz = _top_text(blitz_g, "BLITZ", denom=len(blitz_g)) if len(blitz_g) else "No blitz"
            rows.append([front, front_snaps, _count_pct(front_snaps, total), _count_pct(blitz_count, front_snaps), top_blitz])
        tables = _tables(slide)
        if tables:
            if len(tables[0].table.columns) >= 5:
                _fill_table_preserve(tables[0], ["Front", "Snaps", "Usage", "Blitz %", "Top Blitz Call"], rows)
            else:
                _fill_table_preserve(tables[0], ["Front", "Snaps", "Usage"], [r[:3] for r in rows])
        top_front, n, d, _ = _top_value(df, "FRONT")
        top_third = _top_text(df[df["DOWN"] == 3], "FRONT")
        rz_parts = [g for g in _red_zone_bucket(df).values() if not g.empty]
        rz = pd.concat(rz_parts, ignore_index=True) if rz_parts else pd.DataFrame(columns=df.columns)
        rz_front = _top_text(rz, "FRONT") if not rz.empty else "No data"
        _write_takeaways(slide, [
            f"🔥 Base front: {top_front} — {_count_pct(n, d)}",
            f"⚠️ 3rd down front: {top_third}",
            f"🛑 Red zone front: {rz_front}",
            "AI summary: start with front/blitz relationship."
        ])

    # Blitzes
    if idx["blitz"] is not None:
        slide = prs.slides[idx["blitz"]]
        blitz_df = df[df["IS_BLITZ"] & ~df["BLITZ"].isin(BLANK_BLITZ)].copy()
        blitz_freq = frequency(blitz_df, "BLITZ", denom=total).head(5)
        rows = []
        for _, r in blitz_freq.iterrows():
            b = r["BLITZ"]
            g = blitz_df[blitz_df["BLITZ"] == b]
            blitz_snaps = len(g)
            rows.append([
                b,
                blitz_snaps,
                _count_pct(blitz_snaps, total),
                _top_text(g, "FRONT", denom=blitz_snaps),
                _top_text(g, "STUNT", denom=blitz_snaps),
            ])
        tables = _tables(slide)
        if tables:
            # Keep the existing 5-column table: add top front and embed count/total/% inside top front and top stunt cells.
            _fill_table_preserve(tables[0], ["Blitz", "Snaps", "Usage", "Top Front", "Top Stunt"], rows)
        top_blitz = _top_text(blitz_df, "BLITZ", denom=total) if len(blitz_df) else "No blitz data"
        third_blitz = _top_text(blitz_df[blitz_df["DOWN"] == 3], "BLITZ") if len(blitz_df) else "No blitz data"
        _write_takeaways(slide, [
            f"🔥 Top blitz: {top_blitz}",
            f"⚠️ 3rd down pressure: {third_blitz}",
            f"🛑 Overall blitz rate: {_count_pct(int(df['IS_BLITZ'].sum()), total)}",
            "AI summary: blank blitz cells are excluded."
        ])

    # Coverages
    if idx["coverage"] is not None:
        slide = prs.slides[idx["coverage"]]
        cover_freq = frequency(df, "COVERAGE").head(5)
        rows = []
        for _, r in cover_freq.iterrows():
            cov = r["COVERAGE"]
            cov_count = int((df["COVERAGE"] == cov).sum())
            row = [cov, _count_pct(cov_count, total)]
            for down in [1, 2, 3, 4]:
                gd = df[df["DOWN"] == down]
                row.append(_count_pct(int((gd["COVERAGE"] == cov).sum()), len(gd)))
            rows.append(row)
        tables = _tables(slide)
        if tables:
            _fill_table_preserve(tables[0], ["Coverage", "Overall", "1st Down", "2nd Down", "3rd Down", "4th Down"], rows)
        top_cov = _top_text(df, "COVERAGE")
        man_count = int(df["IS_DISRESPECTFUL"].sum())
        _write_takeaways(slide, [
            f"🔥 Base coverage: {top_cov}",
            f"⚠️ Man/press rate: {_count_pct(man_count, total)}",
            "🛑 Cover 0/Cover 1/press = shot alert.",
            "AI summary: disrespect will not be tolerated."
        ])

    # 3rd down
    if idx["third"] is not None:
        slide = prs.slides[idx["third"]]
        buckets = _third_down_bucket(df)
        tables = _tables(slide)
        if tables:
            headers, rows = _situational_rows_for_template(buckets, tables[0])
            _fill_table_preserve(tables[0], headers, rows)
        all3 = df[df["DOWN"] == 3]
        top_combo = _top_combos_text(all3, 1)
        _write_takeaways(slide, [
            f"🔥 Top 3rd down call: {top_combo}",
            f"⚠️ 3rd down blitz rate: {_count_pct(int(all3['IS_BLITZ'].sum()), len(all3))}",
            f"🛑 3rd down sample: {len(all3)} snaps ({confidence(len(all3))} confidence)",
            "AI summary: protection plan starts here."
        ])

    # Red Zone
    if idx["redzone"] is not None:
        slide = prs.slides[idx["redzone"]]
        buckets = _red_zone_bucket(df)
        tables = _tables(slide)
        if tables:
            headers, rows = _situational_rows_for_template(buckets, tables[0])
            _fill_table_preserve(tables[0], headers, rows)
        rz_parts = [g for g in buckets.values() if not g.empty]
        rz = pd.concat(rz_parts, ignore_index=True) if rz_parts else pd.DataFrame(columns=df.columns)
        _write_takeaways(slide, [
            f"🔥 Top RZ call: {_top_combos_text(rz, 1)}",
            f"⚠️ RZ blitz rate: {_count_pct(int(rz['IS_BLITZ'].sum()) if not rz.empty else 0, len(rz))}",
            f"🛑 Red zone sample: {len(rz)} snaps ({confidence(len(rz))} confidence)",
            "AI summary: finish drives with answers."
        ])

    # Formation
    if idx["formation"] is not None:
        slide = prs.slides[idx["formation"]]
        form_freq = frequency(df, "FORMATION").head(5)
        rows = []
        for _, r in form_freq.iterrows():
            form = r["FORMATION"]
            g = df[df["FORMATION"] == form]
            top_cov = _top_text(g, "COVERAGE", denom=len(g))
            blitz_count = int(g["IS_BLITZ"].sum())
            rows.append([form, _top_fronts_for_formation(g), _count_pct(blitz_count, len(g)), top_cov])
        tables = _tables(slide)
        if tables:
            _fill_table_preserve(tables[0], ["Formation", "Front", "Blitz", "Coverage"], rows)
        top_form = _top_text(df, "FORMATION")
        _write_takeaways(slide, [
            f"🔥 Most frequent formation: {top_form}",
            "⚠️ If top front is under 50%, top two fronts are listed.",
            "🛑 Check RZ formation tendencies before goal-line calls.",
            "AI summary: formation splits are strong predictors."
        ])

    bio = io.BytesIO()
    prs.save(bio)
    return bio.getvalue()
